from io import BytesIO

from pptx import Presentation

from core.color_extractor import ClientConfig, extract_palette
from core.eda import run as run_eda
from core.indicators import derive
from core.metrics import compute
from core.ppt_builder import PPTBuilder


def test_build_pptx_for_fleischmann(excel_bytes, client_logo_bytes, provider_logo_bytes):
    report = run_eda(excel_bytes)
    kpis = compute(report)
    bundle = derive(report, kpis, "Fleischmann")
    palette = extract_palette(client_logo_bytes)
    config = ClientConfig(name="Fleischmann", logo_bytes=client_logo_bytes, palette=palette)

    builder = PPTBuilder(config)
    pptx_bytes = builder.build(kpis, bundle, provider_logo_bytes)

    assert pptx_bytes
    assert len(pptx_bytes) > 50_000  # sanity check

    prs = Presentation(BytesIO(pptx_bytes))
    assert len(prs.slides) == 13
    # Cover should not have charts
    cover = prs.slides[0]
    n_charts_cover = sum(1 for sh in cover.shapes if sh.shape_type == 3)
    assert n_charts_cover == 0
    # At least one chart in the body slides (volume_month onward)
    n_charts_total = sum(
        1 for slide in prs.slides for sh in slide.shapes if sh.shape_type == 3
    )
    assert n_charts_total >= 8
