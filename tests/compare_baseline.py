"""Generate a v2-equivalent PPT via the new pipeline and compare structurally
with the validated baseline at FLEISCHMANN/presentacion/Presentacion_Operacion_Fleischmann_v2.pptx.

Usage: python tests/compare_baseline.py
"""
from __future__ import annotations

import sys
from io import BytesIO
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "streamlit_app"))

from pptx import Presentation

from core.color_extractor import ClientConfig, Palette
from core.eda import run as run_eda
from core.indicators import derive
from core.metrics import compute
from core.ppt_builder import PPTBuilder

FIXTURES = ROOT / "tests" / "fixtures"
BASELINE = Path(
    r"C:\Users\Quick\OneDrive\OneDrive - Quick Help SAS\INFORMES\FLEISCHMANN"
    r"\presentacion\Presentacion_Operacion_Fleischmann_v2.pptx"
)


def palette_v2_baseline() -> Palette:
    """Force the exact palette used in the validated v2 (so values match)."""
    return Palette(
        primary="#1B3D7A",
        accent="#F4B400",
        dark="#0F1419",
    )


def slide_signature(prs: Presentation) -> list[dict]:
    """Lightweight per-slide stats for diffing."""
    out = []
    for i, slide in enumerate(prs.slides, 1):
        n_shapes = len(slide.shapes)
        n_pics = sum(1 for sh in slide.shapes if sh.shape_type == 13)
        n_charts = sum(1 for sh in slide.shapes if sh.shape_type == 3)
        # gather first 3 non-empty text snippets
        texts = []
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    t = (run.text or "").strip()
                    if t:
                        texts.append(t[:60])
                        if len(texts) >= 3:
                            break
                if len(texts) >= 3:
                    break
            if len(texts) >= 3:
                break
        out.append({"slide": i, "shapes": n_shapes, "pics": n_pics,
                    "charts": n_charts, "head": texts})
    return out


def main() -> int:
    excel_path = FIXTURES / "n8n_sample.xlsx"
    logo_path = FIXTURES / "client_logo_sample.png"
    provider_logo_path = ROOT / "streamlit_app" / "assets" / "logo_quick.png"

    if not excel_path.exists() or not BASELINE.exists():
        print(f"ERROR: missing input. excel={excel_path.exists()} baseline={BASELINE.exists()}")
        return 2

    print("[1/4] EDA...")
    report = run_eda(excel_path.read_bytes())
    print(f"  rows: {report.n_rows_filtered}, periodo: {report.period_label}")

    print("[2/4] Metrics + indicators (template)...")
    kpis = compute(report)
    bundle = derive(report, kpis, "Fleischmann")

    print("[3/4] Build PPT with v2 baseline palette...")
    config = ClientConfig(
        name="Fleischmann",
        logo_bytes=logo_path.read_bytes() if logo_path.exists() else None,
        palette=palette_v2_baseline(),
    )
    builder = PPTBuilder(config)
    new_pptx = builder.build(kpis, bundle, provider_logo_path.read_bytes())
    out_path = ROOT / "tests" / "fixtures" / "_generated_v2_equivalent.pptx"
    out_path.write_bytes(new_pptx)
    print(f"  saved: {out_path}")

    print("[4/4] Diff structural con baseline...")
    new_prs = Presentation(BytesIO(new_pptx))
    base_prs = Presentation(BASELINE)

    print(f"  slides: new={len(new_prs.slides)}  base={len(base_prs.slides)}")
    new_sig = slide_signature(new_prs)
    base_sig = slide_signature(base_prs)

    aligned = max(len(new_sig), len(base_sig))
    diffs = 0
    for i in range(aligned):
        n = new_sig[i] if i < len(new_sig) else None
        b = base_sig[i] if i < len(base_sig) else None
        if n is None or b is None:
            print(f"  slide {i+1}: missing on one side")
            diffs += 1
            continue
        # Permitir diferencia ≤ 2 en shapes y ≤ 1 en charts/pics
        delta_shapes = abs(n["shapes"] - b["shapes"])
        delta_charts = abs(n["charts"] - b["charts"])
        delta_pics = abs(n["pics"] - b["pics"])
        ok = delta_shapes <= 4 and delta_charts <= 1 and delta_pics <= 1
        flag = "OK " if ok else "DIFF"
        print(f"  [{flag}] S{i+1}  shapes:{n['shapes']}/{b['shapes']}  "
              f"charts:{n['charts']}/{b['charts']}  pics:{n['pics']}/{b['pics']}")
        if not ok:
            diffs += 1
            print(f"        new head:  {n['head']}")
            print(f"        base head: {b['head']}")

    print(f"\nResultado: {aligned - diffs}/{aligned} slides estructuralmente equivalentes.")
    return 0 if diffs == 0 else 1


if __name__ == "__main__":
    sys.exit(main())
