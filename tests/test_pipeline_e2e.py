"""End-to-end pipeline test: bytes in → bytes out, structurally equivalent
to the Fleischmann v2 baseline."""
from __future__ import annotations

from io import BytesIO

from pptx import Presentation

from core.pipeline import run_pipeline_full


def test_pipeline_full_round_trip(excel_bytes, client_logo_bytes):
    log_lines: list[str] = []
    result = run_pipeline_full(
        excel_bytes,
        client_logo_bytes,
        "Cliente Demo",
        use_llm=False,
        progress_cb=log_lines.append,
    )
    # Output sanity
    assert "pptx_bytes" in result
    assert isinstance(result["pptx_bytes"], (bytes, bytearray))
    assert len(result["pptx_bytes"]) > 50_000
    assert result["used_llm"] is False
    assert result["rows_filtered"] == 2414
    assert result["client_name"] == "Cliente Demo"
    assert result["palette_primary"].startswith("#")
    assert result["n_indicators"] >= 8

    # Progress callback fired
    assert len(log_lines) >= 5

    # PPT has exactly 13 slides
    prs = Presentation(BytesIO(result["pptx_bytes"]))
    assert len(prs.slides) == 13

    # Cover and closing have pictures (logos)
    cover_pics = sum(1 for sh in prs.slides[0].shapes if sh.shape_type == 13)
    closing_pics = sum(1 for sh in prs.slides[12].shapes if sh.shape_type == 13)
    assert cover_pics >= 1
    assert closing_pics >= 1


def test_pipeline_without_logo_uses_default_palette(excel_bytes):
    """When client logo is absent the pipeline still works, with default palette."""
    result = run_pipeline_full(
        excel_bytes,
        None,  # no logo
        "Sin Logo",
        use_llm=False,
    )
    assert result["palette_primary"] == "#1B3D7A"  # default
    prs = Presentation(BytesIO(result["pptx_bytes"]))
    assert len(prs.slides) == 13


def test_pipeline_progress_callback_messages(excel_bytes, client_logo_bytes):
    seen: list[str] = []
    run_pipeline_full(
        excel_bytes, client_logo_bytes, "Cliente Demo",
        use_llm=False, progress_cb=seen.append,
    )
    txt = " · ".join(seen)
    assert "Validando" in txt
    assert "exploratorio" in txt or "Análisis" in txt
    assert "PPT" in txt
    assert "Listo" in txt
