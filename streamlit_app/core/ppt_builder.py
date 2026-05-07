"""Parametrized PPT builder. Same visual identity as the validated v2, but every
brand element (colours, logos, client name) comes from `ClientConfig`.

Usage:
    builder = PPTBuilder(client_config, slides_config)
    pptx_bytes = builder.build(kpis, indicators, provider_logo_bytes)
"""
from __future__ import annotations

from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from .color_extractor import ClientConfig, Palette
from .indicators import IndicatorBundle
from .metrics import Kpis, fmt_int


SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

DEFAULT_SECTIONS = {
    "cover":         (0, "", "primary"),
    "exec_summary":  (1, "1 · RITMO OPERATIVO", "accent"),
    "volume_month":  (1, "1 · RITMO OPERATIVO", "accent"),
    "service_mix":   (1, "1 · RITMO OPERATIVO", "accent"),
    "by_client":     (2, "2 · CLIENTE", "primary"),
    "top_routes":    (2, "2 · CLIENTE", "primary"),
    "top_drivers":   (3, "3 · EQUIPO QUICK HELP", "dark"),
    "vehicle_types": (3, "3 · EQUIPO QUICK HELP", "dark"),
    "keepers":       (3, "3 · EQUIPO QUICK HELP", "dark"),
    "manifests":     (4, "4 · CALIDAD Y TRAZABILIDAD", "success"),
    "cargo_types":   (4, "4 · CALIDAD Y TRAZABILIDAD", "success"),
    "conclusions":   (5, "5 · CONCLUSIONES", "accent"),
    "closing":       (0, "", "primary"),
}


def _hx(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class PPTBuilder:
    """Builds the managerial PPT for any client based on a parametrized palette."""

    def __init__(
        self,
        client: ClientConfig,
        slides_config: dict | None = None,
        *,
        provider_name: str = "Quick Help SAS",
    ):
        self.client = client
        self.palette: Palette = client.palette
        self.slides_config = slides_config
        self.provider_name = provider_name

        # Working state populated during build()
        self.prs: Presentation | None = None
        self.kpis: Kpis | None = None
        self.bundle: IndicatorBundle | None = None
        self.provider_logo_bytes: bytes | None = None
        self._page = 0
        self._total = 0

    # ---------------------------------------------------------------------
    # Public entry
    # ---------------------------------------------------------------------
    def build(
        self,
        kpis: Kpis,
        indicators: IndicatorBundle,
        provider_logo_bytes: bytes,
    ) -> bytes:
        self.kpis = kpis
        self.bundle = indicators
        self.provider_logo_bytes = provider_logo_bytes

        slide_ids = self._effective_slide_order()
        self._total = len(slide_ids)
        self._page = 0

        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H

        for sid in slide_ids:
            handler = getattr(self, f"_slide_{sid}", None)
            if handler is None:
                continue
            handler()

        out = BytesIO()
        self.prs.save(out)
        return out.getvalue()

    # ---------------------------------------------------------------------
    # Slide ordering
    # ---------------------------------------------------------------------
    def _effective_slide_order(self) -> list[str]:
        cfg = self.slides_config
        if not cfg:
            return list(DEFAULT_SECTIONS.keys())
        slides = sorted(
            (s for s in cfg.get("slides", []) if s.get("enabled", True)),
            key=lambda s: s.get("order", 999),
        )
        return [s["id"] for s in slides]

    # ---------------------------------------------------------------------
    # Drawing primitives
    # ---------------------------------------------------------------------
    def _new_slide(self):
        self._page += 1
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _add_rect(self, slide, x, y, w, h, fill, line=None):
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = _hx(fill)
        if line is None:
            sh.line.fill.background()
        else:
            sh.line.color.rgb = _hx(line)
            sh.line.width = Pt(0.5)
        sh.shadow.inherit = False
        return sh

    def _add_text(self, slide, x, y, w, h, text, *, size=14, bold=False, color=None,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
                  italic=False, spacing=1.0):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = align
        p.line_spacing = spacing
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = _hx(color or self.palette.dark)
        run.font.name = font
        return tb

    def _add_multitext(self, slide, x, y, w, h, parts, *, anchor=MSO_ANCHOR.TOP,
                       align=PP_ALIGN.LEFT, spacing=1.15):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        for i, (txt, opt) in enumerate(parts):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.line_spacing = spacing
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(opt.get("size", 12))
            r.font.bold = opt.get("bold", False)
            r.font.italic = opt.get("italic", False)
            r.font.color.rgb = _hx(opt.get("color", self.palette.dark))
            r.font.name = opt.get("font", "Calibri")
        return tb

    def _add_pic(self, slide, image_bytes_or_path, x, y, *, width=None, height=None):
        if image_bytes_or_path is None:
            return None
        try:
            if isinstance(image_bytes_or_path, (bytes, bytearray)):
                img = BytesIO(image_bytes_or_path)
            else:
                img = str(image_bytes_or_path)
            if width:
                return slide.shapes.add_picture(img, x, y, width=width)
            return slide.shapes.add_picture(img, x, y, height=height)
        except Exception:
            return None

    # ---------------------------------------------------------------------
    # Header / footer chrome
    # ---------------------------------------------------------------------
    def _section_for(self, sid: str) -> tuple[int, str, str]:
        return DEFAULT_SECTIONS.get(sid, (0, "", "primary"))

    def _accent_for(self, color_key: str) -> str:
        return getattr(self.palette, color_key, self.palette.accent)

    def _chrome(self, slide, sid: str):
        _, label, _ = self._section_for(sid)
        p = self.palette
        self._add_rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(0.05), p.accent)
        self._add_rect(slide, Emu(0), Inches(0.05), SLIDE_W, Inches(0.45), p.primary)
        self._add_text(slide, Inches(0.35), Inches(0.05), Inches(7), Inches(0.45),
                       label, size=11, bold=True, color=p.white,
                       anchor=MSO_ANCHOR.MIDDLE)
        self._add_text(slide, Inches(8.4), Inches(0.05), Inches(1.4), Inches(0.45),
                       f"{self._page} / {self._total}", size=11, color=p.accent,
                       align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        # Footer
        self._add_rect(slide, Emu(0), Inches(5.32), SLIDE_W, Inches(0.305), p.pale)
        self._add_text(slide, Inches(0.35), Inches(5.33), Inches(9.3), Inches(0.28),
                       f"{self.provider_name}  ·  Operación {self.client.name}  ·  "
                       f"Reporte gerencial {self.kpis['periodo']}",
                       size=9, color=p.grey_dark, anchor=MSO_ANCHOR.MIDDLE)

    def _title(self, slide, kicker: str, title: str, subtitle: str):
        p = self.palette
        self._add_text(slide, Inches(0.35), Inches(0.62), Inches(9.3), Inches(0.32),
                       kicker.upper(), size=10, bold=True, color=p.grey)
        self._add_text(slide, Inches(0.35), Inches(0.92), Inches(9.3), Inches(0.55),
                       title, size=24, bold=True, color=p.primary)
        self._add_text(slide, Inches(0.35), Inches(1.45), Inches(9.3), Inches(0.32),
                       subtitle, size=12, italic=True, color=p.grey_dark)

    # ---------------------------------------------------------------------
    # KPI / insight cards
    # ---------------------------------------------------------------------
    def _kpi_card(self, slide, x, y, w, h, value: str, label: str, accent: str):
        p = self.palette
        self._add_rect(slide, x, y, w, h, p.white, line=p.grey_light)
        self._add_rect(slide, x, y, w, Inches(0.10), accent)
        self._add_text(slide, x + Inches(0.18), y + Inches(0.22),
                       w - Inches(0.36), Inches(0.55),
                       value, size=22, bold=True, color=p.primary)
        self._add_text(slide, x + Inches(0.18), y + Inches(0.70),
                       w - Inches(0.36), Inches(0.32),
                       label, size=9, bold=True, color=p.grey_dark)

    def _kpi_card_dark(self, slide, x, y, w, h, value: str, label: str):
        p = self.palette
        self._add_rect(slide, x, y, w, h, p.dark)
        self._add_rect(slide, x, y, w, Inches(0.10), p.accent)
        self._add_text(slide, x + Inches(0.18), y + Inches(0.22),
                       w - Inches(0.36), Inches(0.55),
                       value, size=22, bold=True, color=p.accent)
        self._add_text(slide, x + Inches(0.18), y + Inches(0.70),
                       w - Inches(0.36), Inches(0.32),
                       label, size=9, bold=True, color=p.white)

    def _insight_card(self, slide, x, y, w, h, kicker: str, big: str, footer: str,
                      *, accent: str | None = None, big_color: str | None = None):
        p = self.palette
        a = accent or p.accent
        bc = big_color or p.primary
        self._add_rect(slide, x, y, w, h, p.white, line=p.grey_light)
        self._add_rect(slide, x, y, Inches(0.12), h, a)
        self._add_text(slide, x + Inches(0.25), y + Inches(0.18),
                       w - Inches(0.4), Inches(0.3),
                       kicker.upper(), size=10, bold=True, color=p.grey_dark)
        self._add_text(slide, x + Inches(0.25), y + Inches(0.50),
                       w - Inches(0.4), Inches(0.7),
                       big, size=24, bold=True, color=bc)
        self._add_text(slide, x + Inches(0.25), y + Inches(1.20),
                       w - Inches(0.4), h - Inches(1.3),
                       footer, size=10, color=p.grey_dark, spacing=1.2)

    def _dark_insight_card(self, slide, x, y, w, h, kicker: str, big: str, footer: str):
        p = self.palette
        self._add_rect(slide, x, y, w, h, p.dark)
        self._add_rect(slide, x, y, Inches(0.12), h, p.accent)
        self._add_text(slide, x + Inches(0.25), y + Inches(0.18),
                       w - Inches(0.4), Inches(0.3),
                       kicker.upper(), size=10, bold=True, color=p.accent)
        self._add_text(slide, x + Inches(0.25), y + Inches(0.50),
                       w - Inches(0.4), Inches(0.7),
                       big, size=26, bold=True, color=p.accent)
        self._add_text(slide, x + Inches(0.25), y + Inches(1.25),
                       w - Inches(0.4), h - Inches(1.35),
                       footer, size=10, color=p.white, spacing=1.2)

    # ---------------------------------------------------------------------
    # Chart helpers
    # ---------------------------------------------------------------------
    def _style_axes(self, chart, num_format=None):
        p = self.palette
        try:
            ca = chart.category_axis
            ca.tick_labels.font.size = Pt(9)
            ca.tick_labels.font.color.rgb = _hx(p.grey_dark)
        except Exception:
            pass
        try:
            va = chart.value_axis
            va.tick_labels.font.size = Pt(9)
            va.tick_labels.font.color.rgb = _hx(p.grey_dark)
            if num_format:
                va.tick_labels.number_format = num_format
                va.tick_labels.number_format_is_linked = False
        except Exception:
            pass

    def _color_series(self, series, color_hex):
        f = series.format.fill
        f.solid()
        f.fore_color.rgb = _hx(color_hex)

    def _chart_column(self, slide, x, y, w, h, cats, vals, color):
        cd = CategoryChartData()
        cd.categories = list(cats)
        cd.add_series("Servicios", list(vals))
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, cd).chart
        chart.has_legend = False
        chart.has_title = False
        plot = chart.plots[0]
        plot.gap_width = 60
        self._color_series(plot.series[0], color)
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.show_value = True
        dl.font.size = Pt(9)
        dl.font.bold = True
        dl.font.color.rgb = _hx(self.palette.dark)
        dl.number_format = "#,##0"
        dl.number_format_is_linked = False
        dl.position = XL_LABEL_POSITION.OUTSIDE_END
        self._style_axes(chart, num_format="#,##0")
        return chart

    def _chart_bar(self, slide, x, y, w, h, cats, vals, color):
        cd = CategoryChartData()
        cd.categories = list(cats)
        cd.add_series("Servicios", list(vals))
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, x, y, w, h, cd).chart
        chart.has_legend = False
        chart.has_title = False
        plot = chart.plots[0]
        plot.gap_width = 60
        self._color_series(plot.series[0], color)
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.show_value = True
        dl.font.size = Pt(9)
        dl.font.bold = True
        dl.font.color.rgb = _hx(self.palette.dark)
        dl.number_format = "#,##0"
        dl.number_format_is_linked = False
        dl.position = XL_LABEL_POSITION.OUTSIDE_END
        self._style_axes(chart, num_format="#,##0")
        return chart

    def _chart_doughnut(self, slide, x, y, w, h, cats, vals, colors):
        cd = CategoryChartData()
        cd.categories = list(cats)
        cd.add_series("", list(vals))
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT, x, y, w, h, cd).chart
        chart.has_title = False
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)
        chart.legend.font.color.rgb = _hx(self.palette.grey_dark)
        plot = chart.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.show_percentage = True
        dl.font.size = Pt(10)
        dl.font.bold = True
        dl.font.color.rgb = _hx(self.palette.white)
        dl.number_format = "0.0%"
        dl.number_format_is_linked = False
        for i, c in enumerate(colors):
            try:
                pt = plot.series[0].points[i]
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = _hx(c)
                pt.format.line.color.rgb = _hx(self.palette.white)
            except Exception:
                pass
        return chart

    def _chart_line(self, slide, x, y, w, h, cats, vals, color):
        cd = CategoryChartData()
        cd.categories = list(cats)
        cd.add_series("Tasa", list(vals))
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, x, y, w, h, cd).chart
        chart.has_legend = False
        chart.has_title = False
        s = chart.plots[0].series[0]
        s.format.line.color.rgb = _hx(color)
        s.format.line.width = Pt(2.5)
        plot = chart.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.show_value = True
        dl.font.size = Pt(9)
        dl.font.bold = True
        dl.font.color.rgb = _hx(self.palette.dark)
        dl.number_format = '0.0"%"'
        dl.number_format_is_linked = False
        dl.position = XL_LABEL_POSITION.ABOVE
        self._style_axes(chart)
        return chart

    # ---------------------------------------------------------------------
    # Slides
    # ---------------------------------------------------------------------
    def _slide_cover(self):
        s = self._new_slide()
        p = self.palette
        k = self.kpis
        self._add_rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, p.primary)
        self._add_rect(s, Emu(0), Inches(3.4), SLIDE_W, Inches(0.07), p.accent)
        self._add_rect(s, Emu(0), Inches(3.47), SLIDE_W, Inches(2.155), p.white)

        self._add_pic(s, self.provider_logo_bytes, Inches(0.4), Inches(0.35),
                      height=Inches(0.95))
        self._add_pic(s, self.client.logo_bytes, Inches(8.0), Inches(3.7),
                      height=Inches(0.85))

        self._add_text(s, Inches(0.4), Inches(1.55), Inches(8), Inches(0.4),
                       "INFORME GERENCIAL DE OPERACIÓN", size=14, bold=True,
                       color=p.accent)
        self._add_text(s, Inches(0.4), Inches(1.95), Inches(9), Inches(0.7),
                       self.client.name, size=44, bold=True, color=p.white)
        self._add_text(s, Inches(0.4), Inches(2.65), Inches(9), Inches(0.4),
                       self.client.subtitle, size=16, color="#D7DEEC")

        cards = [
            (fmt_int(k["total"]), "Servicios ejecutados", p.accent),
            (str(k["n_months"]),  "Meses analizados",     p.primary),
            (str(k["n_clients"]), "Sedes",                p.primary),
            (f"{k['terminal_rate']:.1f}%", "Cumplimiento", p.success),
        ]
        cw = Inches(2.20); ch = Inches(1.30); cx = Inches(0.4); cy = Inches(3.75)
        for v, label, color in cards:
            self._add_rect(s, cx, cy, cw, ch, p.white, line=p.grey_light)
            self._add_rect(s, cx, cy, cw, Inches(0.10), color)
            self._add_text(s, cx + Inches(0.20), cy + Inches(0.22),
                           cw - Inches(0.4), Inches(0.7),
                           v, size=26, bold=True, color=p.primary)
            self._add_text(s, cx + Inches(0.20), cy + Inches(0.85),
                           cw - Inches(0.4), Inches(0.32),
                           label, size=10, bold=True, color=p.grey_dark)
            cx += cw + Inches(0.07)

        self._add_text(s, Inches(0.4), Inches(5.18), Inches(4), Inches(0.18),
                       "PREPARADO POR", size=8, bold=True, color=p.grey)
        self._add_text(s, Inches(0.4), Inches(5.36), Inches(4), Inches(0.22),
                       f"{self.provider_name} · Operador logístico",
                       size=10, bold=True, color=p.dark)
        self._add_text(s, Inches(6.0), Inches(5.18), Inches(3.6), Inches(0.18),
                       "PREPARADO PARA", size=8, bold=True, color=p.grey,
                       align=PP_ALIGN.RIGHT)
        self._add_text(s, Inches(6.0), Inches(5.36), Inches(3.6), Inches(0.22),
                       f"{self.client.name} · Cliente · Confidencial",
                       size=10, bold=True, color=p.primary, align=PP_ALIGN.RIGHT)

    def _slide_exec_summary(self):
        s = self._new_slide()
        p = self.palette
        k = self.kpis
        self._chrome(s, "exec_summary")
        self._title(s, "Resumen ejecutivo",
                    "La operación en una mirada",
                    f"Indicadores clave del período {k['periodo']}")

        grid = [
            (fmt_int(k["total"]),      "Servicios ejecutados",  p.accent),
            (str(k["n_clients"]),      "Sedes",                  p.primary),
            (str(k["n_routes"]),       "Rutas atendidas",        p.primary),
            (str(k["n_drivers"]),      "Conductores activos",    p.accent),
            (str(k["n_keepers"]),      "Tenedores de flota",     p.primary),
            (str(k["n_vehicle_types"]), "Tipos de vehículo",     p.accent),
            (str(k["n_months"]),       "Meses cubiertos",        p.accent),
            (f"{k['terminal_rate']:.1f}%", "Cumplimiento",       p.success),
        ]
        cw, ch = Inches(2.20), Inches(1.05)
        sx, sy = Inches(0.40), Inches(1.95)
        for i, (v, label, color) in enumerate(grid):
            x = sx + (i % 4) * (cw + Inches(0.10))
            y = sy + (i // 4) * (ch + Inches(0.10))
            if i == 7:
                self._kpi_card_dark(s, x, y, cw, ch, v, label)
            else:
                self._kpi_card(s, x, y, cw, ch, v, label, color)

        by = Inches(4.30); bh = Inches(0.95)
        self._add_rect(s, Inches(0.40), by, Inches(9.20), bh, p.pale)
        self._add_rect(s, Inches(0.40), by, Inches(0.10), bh, p.primary)
        self._add_text(s, Inches(0.65), by + Inches(0.10),
                       Inches(8.9), Inches(0.30),
                       "LECTURA EJECUTIVA", size=10, bold=True, color=p.primary)
        self._add_text(s, Inches(0.65), by + Inches(0.38),
                       Inches(8.9), Inches(0.55),
                       self.bundle.executive_summary,
                       size=11, color=p.grey_dark, spacing=1.2)

    def _slide_volume_month(self):
        s = self._new_slide()
        p = self.palette
        k = self.kpis
        self._chrome(s, "volume_month")
        self._title(s, "Sección 1 · Ritmo operativo",
                    "Volumen mensual de servicios",
                    "El pulso de la operación: una curva con estacionalidad clara")

        self._chart_column(s, Inches(0.35), Inches(1.85), Inches(6.0), Inches(3.3),
                           k["vol_categories"], k["vol_values"], p.accent)
        self._dark_insight_card(s, Inches(6.55), Inches(1.85), Inches(3.10), Inches(1.55),
                                "Mes pico", k["peak_month"],
                                f"{fmt_int(k['peak_val'])} servicios · {k['peak_x']:.1f}× sobre el promedio")
        self._insight_card(s, Inches(6.55), Inches(3.55), Inches(3.10), Inches(1.60),
                           "Mes valle", k["valley_month"],
                           f"{fmt_int(k['valley_val'])} servicios · inicio de operación / rampa",
                           accent=p.primary, big_color=p.primary)

    def _slide_service_mix(self):
        s = self._new_slide()
        p = self.palette
        k = self.kpis
        self._chrome(s, "service_mix")
        self._title(s, "Sección 1 · Ritmo operativo",
                    "Mix de tipo de servicio",
                    f"{k['st_top']} lidera la modalidad")

        donut_colors = [p.accent, p.primary, p.success, p.alert, p.dark][:len(k["st_cats"])]
        self._chart_doughnut(s, Inches(0.35), Inches(1.85), Inches(5.4), Inches(3.30),
                             k["st_cats"], k["st_vals"], donut_colors)

        cx, cy = Inches(6.05), Inches(1.85)
        cw, ch = Inches(3.60), Inches(1.05)
        self._insight_card(s, cx, cy, cw, ch, "Modalidad dominante",
                           f"{k['st_top_pct']:.1f}%",
                           f"{k['st_top']} concentra el {k['st_top_pct']:.1f}% del volumen total.",
                           accent=p.accent)
        self._insight_card(s, cx, cy + Inches(1.15), cw, ch, "Diversidad",
                           f"{k['n_service_types']} tipos",
                           "Mix relativamente concentrado, fácil de planear capacidad.",
                           accent=p.primary)
        self._dark_insight_card(s, cx, cy + Inches(2.30), cw, ch,
                                "Recomendación", "Monitoreo mensual",
                                "Revisar desviaciones del mix mes a mes para anticipar capacidad.")

    def _slide_by_client(self):
        s = self._new_slide()
        p = self.palette
        k = self.kpis
        self._chrome(s, "by_client")
        self._title(s, "Sección 2 · Cliente",
                    f"Distribución por sede {self.client.name}",
                    f"{k['top_cli']} concentra el mayor volumen — eje de planeación")

        self._chart_bar(s, Inches(0.35), Inches(1.85), Inches(6.0), Inches(3.3),
                        k["cli_cats"], k["cli_vals"], p.primary)

        cx, cy = Inches(6.55), Inches(1.85); cw = Inches(3.10)
        self._insight_card(s, cx, cy, cw, Inches(1.55), "Sede líder",
                           f"{k['top_cli_pct']:.1f}%",
                           f"{k['top_cli']} concentra esta porción del volumen.",
                           accent=p.primary, big_color=p.primary)
        self._dark_insight_card(s, cx, cy + Inches(1.65), cw, Inches(1.65),
                                "Recomendación", "Anclar capacidad",
                                f"Alinear planeación con {k['top_cli']} (sede ancla).")

    def _slide_top_routes(self):
        s = self._new_slide()
        p = self.palette
        k = self.kpis
        self._chrome(s, "top_routes")
        self._title(s, "Sección 2 · Cliente",
                    "Top rutas origen → destino",
                    f"La ruta líder concentra el {k['top_route_pct']:.1f}% del volumen")

        self._chart_bar(s, Inches(0.35), Inches(1.85), Inches(6.0), Inches(3.3),
                        k["route_cats"], k["route_vals"], p.accent)

        cx = Inches(6.55); cw = Inches(3.10)
        self._insight_card(s, cx, Inches(1.85), cw, Inches(1.00),
                           "Rutas únicas", str(k["n_routes"]),
                           "Cobertura de pares origen-destino atendidos.",
                           accent=p.primary, big_color=p.primary)
        self._insight_card(s, cx, Inches(2.95), cw, Inches(1.00),
                           "Ruta líder", f"{k['top_route_pct']:.1f}%",
                           f"En la ruta {k['top_route']}.",
                           accent=p.accent, big_color=p.accent)
        self._dark_insight_card(s, cx, Inches(4.05), cw, Inches(1.10),
                                "Concentración", f"{k['top10_route_pct']:.1f}%",
                                "del volumen está concentrado en el Top 10 de rutas.")

    def _slide_top_drivers(self):
        s = self._new_slide()
        p = self.palette
        k = self.kpis
        self._chrome(s, "top_drivers")
        self._title(s, "Sección 3 · Equipo",
                    "Top 10 conductores",
                    f"Concentración: 10 conductores hacen el {k['top10_drv_pct']:.1f}% de la operación")

        self._chart_bar(s, Inches(0.35), Inches(1.85), Inches(6.0), Inches(3.3),
                        k["drv_cats"], k["drv_vals"], p.dark)

        cx, cy = Inches(6.55), Inches(1.85); cw = Inches(3.10)
        self._dark_insight_card(s, cx, cy, cw, Inches(1.55),
                                "Conductor líder", k["top_drv"],
                                f"{k['top_drv_pct']:.1f}% del volumen total.")
        self._insight_card(s, cx, cy + Inches(1.65), cw, Inches(1.65),
                           "Concentración", f"{k['top10_drv_pct']:.1f}%",
                           "del volumen lo ejecutan 10 conductores. Plan de retención y backup.",
                           accent=p.alert, big_color=p.alert)

    def _slide_vehicle_types(self):
        s = self._new_slide()
        p = self.palette
        k = self.kpis
        self._chrome(s, "vehicle_types")
        self._title(s, "Sección 3 · Equipo",
                    "Distribución por tipo de vehículo",
                    f"{k['top_veh']} ancla la flota")

        self._chart_column(s, Inches(0.35), Inches(1.85), Inches(6.0), Inches(3.3),
                           k["veh_cats"], k["veh_vals"], p.accent)

        cx, cy = Inches(6.55), Inches(1.85); cw = Inches(3.10)
        self._insight_card(s, cx, cy, cw, Inches(1.55), "Vehículo dominante",
                           f"{k['top_veh_pct']:.1f}%", k["top_veh"],
                           accent=p.accent, big_color=p.accent)
        self._insight_card(s, cx, cy + Inches(1.65), cw, Inches(1.65),
                           "Lectura", f"{k['n_vehicle_types']} tipos",
                           "La mezcla refleja el mix dominante. Validar asignación tipo-vehículo vs. peso.",
                           accent=p.primary, big_color=p.primary)

    def _slide_keepers(self):
        s = self._new_slide()
        p = self.palette
        k = self.kpis
        self._chrome(s, "keepers")
        self._title(s, "Sección 3 · Equipo",
                    "Concentración por tenedor de flota",
                    f"Pareto: el Top 5 aporta el {k['top5_kee_pct']:.1f}% — riesgo a mitigar")

        self._chart_column(s, Inches(0.35), Inches(1.85), Inches(6.0), Inches(3.3),
                           k["kee_cats"], k["kee_vals"], p.primary)

        cx, cy = Inches(6.55), Inches(1.85); cw = Inches(3.10)
        self._insight_card(s, cx, cy, cw, Inches(1.00),
                           "Tenedores activos", str(k["n_keepers"]),
                           "Aportan vehículos a la operación.",
                           accent=p.primary, big_color=p.primary)
        self._dark_insight_card(s, cx, cy + Inches(1.10), cw, Inches(1.05),
                                "Tenedor líder", f"{k['top_kee_pct']:.1f}%",
                                f"{k['top_kee']} concentra esta porción.")
        self._insight_card(s, cx, cy + Inches(2.25), cw, Inches(1.05),
                           "Pareto Top 5", f"{k['top5_kee_pct']:.1f}%",
                           "del volumen total. Formalizar respaldos contractuales.",
                           accent=p.alert, big_color=p.alert)

    def _slide_manifests(self):
        s = self._new_slide()
        p = self.palette
        k = self.kpis
        self._chrome(s, "manifests")
        self._title(s, "Sección 4 · Calidad y trazabilidad",
                    "Cumplimiento de manifiestos",
                    f"{k['terminal_rate']:.1f}% terminal — operación trazable")

        self._add_rect(s, Inches(0.35), Inches(1.85), Inches(2.95), Inches(1.55), p.dark)
        self._add_rect(s, Inches(0.35), Inches(1.85), Inches(0.12), Inches(1.55), p.success)
        self._add_text(s, Inches(0.55), Inches(1.95), Inches(2.65), Inches(0.30),
                       "TASA TERMINAL", size=10, bold=True, color=p.success_light)
        self._add_text(s, Inches(0.55), Inches(2.18), Inches(2.65), Inches(0.7),
                       f"{k['terminal_rate']:.1f}%", size=42, bold=True, color=p.accent)
        self._add_text(s, Inches(0.55), Inches(2.92), Inches(2.65), Inches(0.45),
                       "Cumplido + Liquidado y Cerrado",
                       size=10, color=p.white, spacing=1.15)

        state_colors = []
        for cat in k["mfs_cats"]:
            state_colors.append({
                "Cumplido": p.success,
                "Liquidado y Cerrado": p.success_light,
                "Activo": p.accent,
                "Anulado": p.alert,
            }.get(cat, p.grey))
        self._chart_doughnut(s, Inches(3.45), Inches(1.65), Inches(3.3), Inches(2.8),
                             k["mfs_cats"], k["mfs_vals"], state_colors)

        lx = Inches(6.85); ly = Inches(1.85); rh = Inches(0.35)
        total_mfs = sum(k["mfs_vals"])
        for cat, val, color in zip(k["mfs_cats"], k["mfs_vals"], state_colors):
            self._add_rect(s, lx, ly, Inches(0.20), rh, color)
            self._add_text(s, lx + Inches(0.30), ly, Inches(2.6), rh,
                           cat, size=10, bold=True, color=p.primary,
                           anchor=MSO_ANCHOR.MIDDLE)
            pct = val / total_mfs * 100
            self._add_text(s, lx + Inches(0.30), ly, Inches(2.6), rh,
                           f"{pct:.1f}%  ·  {fmt_int(val)}",
                           size=10, color=p.grey_dark,
                           anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
            ly += rh + Inches(0.07)

        self._add_text(s, Inches(0.35), Inches(3.55), Inches(9.3), Inches(0.30),
                       "EVOLUCIÓN MENSUAL · TASA TERMINAL",
                       size=10, bold=True, color=p.primary)
        self._chart_line(s, Inches(0.35), Inches(3.85), Inches(9.30), Inches(1.40),
                         k["term_months"], k["term_rates"], p.success)

    def _slide_cargo_types(self):
        s = self._new_slide()
        p = self.palette
        k = self.kpis
        self._chrome(s, "cargo_types")
        self._title(s, "Sección 4 · Calidad y trazabilidad",
                    "Tipos de mercancía transportada",
                    f"{k['top_merc']} concentra el {k['top_merc_pct']:.1f}%")

        self._chart_bar(s, Inches(0.35), Inches(1.85), Inches(6.0), Inches(3.3),
                        k["merc_cats"], k["merc_vals"], p.primary)

        cx, cy = Inches(6.55), Inches(1.85); cw = Inches(3.10)
        self._insight_card(s, cx, cy, cw, Inches(1.55), "Categoría dominante",
                           f"{k['top_merc_pct']:.1f}%", k["top_merc"][:65],
                           accent=p.primary, big_color=p.primary)
        self._insight_card(s, cx, cy + Inches(1.65), cw, Inches(1.65),
                           "Calidad del dato", f"{k['n_merc']} categorías",
                           "Hay duplicados por tildes y mayúsculas. Recomendación: estandarizar el catálogo.",
                           accent=p.accent, big_color=p.accent)

    def _slide_conclusions(self):
        s = self._new_slide()
        p = self.palette
        b = self.bundle
        self._chrome(s, "conclusions")
        self._title(s, "Sección 5 · Conclusiones",
                    "Fortalezas y oportunidades",
                    "Foco para los próximos ciclos operativos")

        col_w = Inches(4.55)
        col_x_left = Inches(0.40); col_x_right = Inches(5.05)
        col_y = Inches(1.85)

        self._add_rect(s, col_x_left, col_y, col_w, Inches(0.45), p.success)
        self._add_text(s, col_x_left + Inches(0.2), col_y, col_w - Inches(0.4), Inches(0.45),
                       "FORTALEZAS DETECTADAS", size=11, bold=True, color=p.white,
                       anchor=MSO_ANCHOR.MIDDLE)
        fy = col_y + Inches(0.55)
        for title, text in b.fortalezas:
            self._add_rect(s, col_x_left, fy, Inches(0.30), Inches(0.50), p.success)
            self._add_text(s, col_x_left, fy, Inches(0.30), Inches(0.50),
                           "✓", size=14, bold=True, color=p.white,
                           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            self._add_multitext(s, col_x_left + Inches(0.40), fy - Inches(0.02),
                                col_w - Inches(0.45), Inches(0.55),
                                [(title + " — ", {"size": 11, "bold": True, "color": p.primary}),
                                 (text, {"size": 10, "color": p.grey_dark})],
                                anchor=MSO_ANCHOR.TOP, spacing=1.15)
            fy += Inches(0.55)

        self._add_rect(s, col_x_right, col_y, col_w, Inches(0.45), p.accent)
        self._add_text(s, col_x_right + Inches(0.2), col_y, col_w - Inches(0.4), Inches(0.45),
                       "OPORTUNIDADES DE MEJORA", size=11, bold=True, color=p.dark,
                       anchor=MSO_ANCHOR.MIDDLE)
        fy = col_y + Inches(0.55)
        for n, (title, text) in enumerate(b.oportunidades, start=1):
            self._add_rect(s, col_x_right, fy, Inches(0.30), Inches(0.50), p.accent)
            self._add_text(s, col_x_right, fy, Inches(0.30), Inches(0.50),
                           str(n), size=14, bold=True, color=p.dark,
                           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            self._add_multitext(s, col_x_right + Inches(0.40), fy - Inches(0.02),
                                col_w - Inches(0.45), Inches(0.55),
                                [(title + " — ", {"size": 11, "bold": True, "color": p.primary}),
                                 (text, {"size": 10, "color": p.grey_dark})],
                                anchor=MSO_ANCHOR.TOP, spacing=1.15)
            fy += Inches(0.55)

    def _slide_closing(self):
        s = self._new_slide()
        p = self.palette
        k = self.kpis
        self._add_rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, p.primary)
        self._add_rect(s, Emu(0), Inches(2.7), SLIDE_W, Inches(0.07), p.accent)

        self._add_text(s, Inches(0.4), Inches(1.0), Inches(9.2), Inches(1.4),
                       "Gracias", size=72, bold=True, color=p.white,
                       align=PP_ALIGN.CENTER)
        self._add_text(s, Inches(0.4), Inches(2.95), Inches(9.2), Inches(0.5),
                       f"Soluciones logísticas para {self.client.name}",
                       size=18, color=p.accent, align=PP_ALIGN.CENTER)

        self._add_pic(s, self.provider_logo_bytes, Inches(1.35), Inches(3.85),
                      height=Inches(0.85))
        self._add_pic(s, self.client.logo_bytes, Inches(7.0), Inches(3.95),
                      height=Inches(0.65))

        self._add_text(s, Inches(0.4), Inches(4.85), Inches(4), Inches(0.30),
                       self.provider_name.upper(), size=11, bold=True, color=p.accent,
                       align=PP_ALIGN.CENTER)
        self._add_text(s, Inches(0.4), Inches(5.10), Inches(4), Inches(0.25),
                       "Equipo de Operaciones y Analítica", size=10, color=p.white,
                       align=PP_ALIGN.CENTER)
        self._add_text(s, Inches(5.6), Inches(4.85), Inches(4), Inches(0.30),
                       self.client.name.upper(), size=11, bold=True, color=p.accent,
                       align=PP_ALIGN.CENTER)
        self._add_text(s, Inches(5.6), Inches(5.10), Inches(4), Inches(0.25),
                       f"Confidencial · {k['periodo']}", size=10, color=p.white,
                       align=PP_ALIGN.CENTER)
